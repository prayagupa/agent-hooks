#!/usr/bin/env python3
"""Build agent-hooks design-review deck.

Layout: 13.333" x 7.5" widescreen. Microsoft brand palette (matches
AAIF_AGT_Technical_Review.pptx). All diagrams drawn as native PowerPoint
shapes so they remain editable and crisp.

Run:  .venv/bin/python build_deck.py
Out:  agent-hooks-design-review.pptx
"""
from __future__ import annotations

import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---- palette (matches AGT deck) -------------------------------------------
INK = RGBColor(0x10, 0x10, 0x10)
GRAPHITE = RGBColor(0x2B, 0x2B, 0x2B)
SLATE = RGBColor(0x4A, 0x4A, 0x55)
SUBTLE = RGBColor(0x6E, 0x6E, 0x6E)
RULE = RGBColor(0xCC, 0xCC, 0xCC)
BG = RGBColor(0xF7, 0xF7, 0xF9)
CARD = RGBColor(0xFF, 0xFF, 0xFF)

BLUE = RGBColor(0x00, 0x78, 0xD4)
CYAN = RGBColor(0x00, 0xB7, 0xC3)
MAGENTA = RGBColor(0xB4, 0x00, 0x9E)
GREEN = RGBColor(0x10, 0x88, 0x3E)
AMBER = RGBColor(0xCA, 0x5A, 0x00)
RED = RGBColor(0xC5, 0x0F, 0x1F)
NAVY = RGBColor(0x10, 0x2A, 0x43)

CODE_BG = RGBColor(0x1B, 0x1F, 0x27)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)

FONT = "Segoe UI"
MONO = "Cascadia Mono"

W, H = 13.333, 7.5

# ---- primitives ------------------------------------------------------------


def txt(tf, text, *, size=14, bold=False, color=GRAPHITE, font=FONT,
        align=None, anchor=None, italic=False, first=True, space_after=2):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    if anchor is not None:
        tf.vertical_anchor = anchor
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def box(slide, x, y, w, h, *, fill=None, line=None, line_w=1.0,
        radius=0.08, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE and radius is not None:
        try:
            sp.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def label(slide, x, y, w, h, text, *, size=12, bold=False, color=GRAPHITE,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT,
          fill=None, line=None, line_w=1.0):
    sp = box(slide, x, y, w, h, fill=fill, line=line, line_w=line_w)
    tf = sp.text_frame
    tf.word_wrap = True
    txt(tf, text, size=size, bold=bold, color=color, align=align,
        anchor=anchor, font=font)
    return sp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def connector(slide, x1, y1, x2, y2, *, color=SLATE, width=1.5,
              dashed=False, arrow=True):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln()
        pd = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(pd)
    if arrow:
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln()
        he = ln.makeelement(qn("a:headEnd"), {"type": "none"})
        te = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(he)
        ln.append(te)
    return c


def header(slide, eyebrow, title, subtitle=None):
    box(slide, 0, 0, W, H, fill=BG, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(slide, 0.7, 0.35, W - 1.4, 1.4)
    txt(tf, eyebrow, size=11, bold=True, color=BLUE)
    txt(tf, title, size=26, bold=True, color=INK, first=False, space_after=4)
    if subtitle:
        txt(tf, subtitle, size=13, color=SLATE, first=False)
    box(slide, 0.7, 1.65, W - 1.4, 0.015, fill=RULE, shape=MSO_SHAPE.RECTANGLE)


def footer(slide, n, total):
    tf = textbox(slide, W - 1.5, H - 0.4, 1.0, 0.3)
    txt(tf, f"{n} / {total}", size=9, color=SUBTLE, align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---- deck ------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]
TOTAL = 14


# ---------------------------------------------------------------------------
# 1. Title
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, H - 0.18, W, 0.18, fill=CYAN, shape=MSO_SHAPE.RECTANGLE)
tf = textbox(s, 1.0, 2.2, W - 2.0, 3.0)
txt(tf, "AGENT-HOOKS", size=13, bold=True, color=CYAN)
txt(tf, "A framework-neutral control contract for AI agents",
    size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), first=False, space_after=10)
txt(tf, "Design review · spec v0.1.0-alpha · responsibleai/agent-hooks",
    size=13, color=RGBColor(0xC8, 0xC8, 0xC8), first=False)
notes(s, "Extracted from ACS policy-engine v0.3.1-beta. Goal: every agent "
         "framework exposes the same eight interception points; every "
         "control component (policy engine, content filter, rate limiter, "
         "approval gateway, egress guard) targets one contract. Control "
         "plane only; tracing and passive observability are out of scope.")


# ---------------------------------------------------------------------------
# 2. The problem: N frameworks x M interceptors
s = prs.slides.add_slide(BLANK)
header(s, "PROBLEM", "Every framework x every control = bespoke adapter",
       "Surveyed 7 framework integrations in AGT: no shared base, divergent "
       "payloads, uneven interception-point coverage")

fws = ["OpenAI Agents", "LangChain", "Semantic Kernel", "CrewAI",
       "PydanticAI", "AutoGen", "LangGraph"]
cons = ["ACS policy", "Content filter", "Rate limiter", "Approval gateway",
        "Egress guard"]

fy, cy = 2.2, 5.2
fw_w, gap = 1.55, 0.13
fx0 = 0.7 + (W - 1.4 - len(fws) * fw_w - (len(fws) - 1) * gap) / 2
for i, f in enumerate(fws):
    label(s, fx0 + i * (fw_w + gap), fy, fw_w, 0.75, f,
          fill=CARD, line=BLUE, size=10, bold=True)

c_w = 2.05
cx0 = 0.7 + (W - 1.4 - len(cons) * c_w - (len(cons) - 1) * gap) / 2
for j, c in enumerate(cons):
    label(s, cx0 + j * (c_w + gap), cy, c_w, 0.75, c,
          fill=CARD, line=MAGENTA, size=10, bold=True)

# tangle: every framework to every interceptor
for i in range(len(fws)):
    for j in range(len(cons)):
        connector(s, fx0 + i * (fw_w + gap) + fw_w / 2, fy + 0.75,
                  cx0 + j * (c_w + gap) + c_w / 2, cy,
                  color=RULE, width=0.75, arrow=False)

label(s, W / 2 - 1.6, 3.85, 3.2, 0.55,
      f"{len(fws)} x {len(cons)} = {len(fws)*len(cons)} adapters",
      fill=AMBER, size=13, bold=True, color=CARD)

footer(s, 2, TOTAL)
notes(s, "agentmesh-integrations survey: payload representation ranges from "
         "typed dataclasses to Option<String>; agent identity resolution is "
         "ad-hoc; PydanticAI/CrewAI lack agent-level lifecycle entirely.")


# ---------------------------------------------------------------------------
# 3. The cut: agent-hooks as the contract layer
s = prs.slides.add_slide(BLANK)
header(s, "ARCHITECTURE", "agent-hooks is the wire contract; ACS becomes one interceptor",
       "Frameworks emit AgentContext; interceptors return Verdict; host enforces")

# top: frameworks
for i, f in enumerate(fws):
    label(s, fx0 + i * (fw_w + gap), 2.0, fw_w, 0.65, f,
          fill=CARD, line=BLUE, size=9, bold=True)
    connector(s, fx0 + i * (fw_w + gap) + fw_w / 2, 2.65,
              fx0 + i * (fw_w + gap) + fw_w / 2, 3.15,
              color=BLUE, width=1.0)

# middle: agent-hooks bar
ah = box(s, 0.9, 3.15, W - 1.8, 1.0, fill=NAVY, line=NAVY)
tf = ah.text_frame
tf.word_wrap = True
txt(tf, "agent-hooks", size=16, bold=True, color=CYAN,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(tf, "InterceptionPoint  ·  AgentContext  ·  Verdict  ·  host obligations  ·  "
        "approval seam  ·  context_identity  ·  CTK",
    size=10, color=RGBColor(0xC8, 0xD4, 0xE0), align=PP_ALIGN.CENTER, first=False)

# down arrows
label(s, 1.5, 4.25, 2.0, 0.3, "AgentContext ↓", size=9, color=SLATE, line=None)
label(s, W - 3.5, 4.25, 2.0, 0.3, "↑ Verdict", size=9, color=SLATE, line=None)

# bottom: interceptors
for j, c in enumerate(cons):
    label(s, cx0 + j * (c_w + gap), 5.0, c_w, 0.65, c,
          fill=CARD, line=MAGENTA, size=10, bold=True)
    connector(s, cx0 + j * (c_w + gap) + c_w / 2, 4.15,
              cx0 + j * (c_w + gap) + c_w / 2, 5.0,
              color=MAGENTA, width=1.0)

label(s, W / 2 - 2.0, 6.05, 4.0, 0.5,
      f"{len(fws)} + {len(cons)} = {len(fws)+len(cons)} adapters",
      fill=GREEN, size=13, bold=True, color=CARD)

footer(s, 3, TOTAL)
notes(s, "ACS keeps PolicyInput, manifest, Cedar/OPA dispatchers, "
         "annotators, runtime_error:* reasons. agent-hooks owns only the "
         "bidirectional wire contract. N+M replaces N*M.")


# ---------------------------------------------------------------------------
# 4. Interception points on the agent-loop timeline
s = prs.slides.add_slide(BLANK)
header(s, "SPEC §3", "Eight interception points on the agent loop",
       "Closed set; pre/post pairs bracket each side-effect; sequence is "
       "strictly increasing")

ty = 3.7
box(s, 0.9, ty - 0.01, W - 1.8, 0.02, fill=SLATE, shape=MSO_SHAPE.RECTANGLE)

points = [
    ("agent_startup", BLUE, False),
    ("input", CYAN, True),
    ("pre_model_call", MAGENTA, True),
    ("post_model_call", MAGENTA, True),
    ("pre_tool_call", AMBER, True),
    ("post_tool_call", AMBER, True),
    ("output", CYAN, True),
    ("agent_shutdown", BLUE, False),
]
xs = [0.9 + i * (W - 1.8) / (len(points) - 1) for i in range(len(points))]

for i, ((name, col, tform), x) in enumerate(zip(points, xs)):
    box(s, x - 0.06, ty - 0.06, 0.12, 0.12, fill=col,
        shape=MSO_SHAPE.OVAL, line=None)
    above = i % 2 == 0
    ly = ty - 1.35 if above else ty + 0.35
    lb = label(s, x - 0.85, ly, 1.7, 0.85, name,
               fill=CARD, line=col, size=10, bold=True, font=MONO)
    if not tform:
        txt(lb.text_frame, "transform forbidden", size=7, color=SUBTLE,
            align=PP_ALIGN.CENTER, first=False)
    connector(s, x, ly + (0.85 if above else 0), x, ty,
              color=col, width=1.0, arrow=False)

# pre/post brackets
for (a, b, lab, col) in [(2, 3, "model call", MAGENTA), (4, 5, "tool call", AMBER)]:
    bx, bw = xs[a], xs[b] - xs[a]
    box(s, bx, ty + 1.5, bw, 0.04, fill=col, shape=MSO_SHAPE.RECTANGLE)
    label(s, bx, ty + 1.6, bw, 0.3, lab, size=9, color=col)

label(s, 0.9, 6.4, W - 1.8, 0.4,
      "§3.1: startup precedes all; shutdown follows all; each pre_X is "
      "followed by exactly one post_X for the same id unless blocked",
      size=10, color=SLATE, fill=None)

footer(s, 4, TOTAL)
notes(s, "transform is forbidden at startup/shutdown (no actionable target). "
         "Capability subsetting (§3.2): a tool-router host with no model "
         "calls may omit *_model_call and still claim conformance.")


# ---------------------------------------------------------------------------
# 5. AgentContext tiers
s = prs.slides.add_slide(BLANK)
header(s, "SPEC §4", "AgentContext: tiered schema with an explicit $target",
       "L0+L1 are sufficient to decide and to compute identity; L2/L3 are "
       "additive")

# nested rectangles
cx, cw = 1.0, 7.4
levels = [
    ("L3", "extensions.<ns>", "namespaced free-for-all (acs, ctk reserved)",
     RGBColor(0xEE, 0xEE, 0xF1), 0.0),
    ("L2", "well-known optional", "trace, tenant, budgets, usage, model.params",
     RGBColor(0xE3, 0xEC, 0xF6), 0.35),
    ("L1", "per-point required", "tool_call.{id,name,args}, response, messages, ...",
     RGBColor(0xD4, 0xE6, 0xF6), 0.7),
    ("L0", "always required", "spec, interception_point, timestamp, sequence, "
     "agent.{id,framework}, session.id, target",
     RGBColor(0xC0, 0xDB, 0xF2), 1.05),
]
top, bot = 2.1, 6.6
for tag, name, body, fill, inset in levels:
    sp = box(s, cx + inset, top + inset, cw - 2 * inset, (bot - top) - 2 * inset,
             fill=fill, line=BLUE, line_w=0.75)
    tf = textbox(s, cx + inset + 0.15, top + inset + 0.08, 2.6, 0.6)
    txt(tf, f"{tag}  {name}", size=11, bold=True, color=NAVY)
    txt(tf, body, size=8, color=SLATE, first=False)

# target callout
tx = cx + cw + 0.5
label(s, tx, 2.4, 3.2, 0.5, "$target", fill=NAVY, size=14, bold=True,
      color=CYAN, font=MONO)
tf = textbox(s, tx, 3.0, 3.2, 3.2)
txt(tf, "The value the action will consume / has produced.",
    size=11, bold=True, color=GRAPHITE)
for hp, tgt in [("input", "input"), ("pre_model_call", "messages"),
                ("post_model_call", "response"), ("pre_tool_call",
                "tool_call.args"), ("post_tool_call", "tool_result.value"),
                ("output", "output")]:
    txt(tf, f"  {hp}", size=9, color=SLATE, font=MONO, first=False, space_after=0)
    txt(tf, f"    -> {tgt}", size=9, color=BLUE, font=MONO, first=False)
txt(tf, "Only $target may be rewritten by transform (§5.2).",
    size=9, color=AMBER, italic=True, first=False)

footer(s, 5, TOTAL)
notes(s, "context_identity = sha256(canonical_json(L0+L1 only)). L2 and "
         "extensions are stripped so adding a trace_id or budget never "
         "changes the approval-binding hash.")


# ---------------------------------------------------------------------------
# 6. Verdict decision tree
s = prs.slides.add_slide(BLANK)
header(s, "SPEC §5–6", "Verdict and host obligations",
       "Five decisions; permit-class proceeds, block-class halts; the host "
       "MUST enforce")

root_x, root_y = W / 2, 2.2
label(s, root_x - 1.0, root_y, 2.0, 0.5, "Verdict.decision",
      fill=NAVY, color=CARD, size=12, bold=True, font=MONO)

decisions = [
    ("allow", GREEN, "proceed; target unchanged", True),
    ("warn", GREEN, "proceed; record reason", True),
    ("transform", CYAN, "apply $target rewrite, then proceed", True),
    ("escalate", AMBER, "halt; consult approval resolver (§9)", False),
    ("deny", RED, "halt; action MUST NOT execute", False),
]
dy, dw = 3.4, 2.25
dx0 = (W - len(decisions) * dw - (len(decisions) - 1) * 0.18) / 2
for i, (name, col, desc, permit) in enumerate(decisions):
    x = dx0 + i * (dw + 0.18)
    connector(s, root_x, root_y + 0.5, x + dw / 2, dy, color=SLATE, width=1.0)
    sp = label(s, x, dy, dw, 0.55, name, fill=col, color=CARD,
               size=12, bold=True, font=MONO)
    label(s, x, dy + 0.6, dw, 0.7, desc, fill=CARD, line=col, size=9)
    connector(s, x + dw / 2, dy + 1.3, x + dw / 2, 5.3,
              color=GREEN if permit else RED, width=1.5)

box(s, dx0, 5.3, 3 * dw + 2 * 0.18, 0.6, fill=None, line=GREEN, line_w=2.0)
label(s, dx0, 5.35, 3 * dw + 2 * 0.18, 0.5, "PERMIT  →  action executes",
      size=11, bold=True, color=GREEN)
box(s, dx0 + 3 * (dw + 0.18), 5.3, 2 * dw + 0.18, 0.6, fill=None,
    line=RED, line_w=2.0)
label(s, dx0 + 3 * (dw + 0.18), 5.35, 2 * dw + 0.18, 0.5,
      "BLOCK  →  action MUST NOT execute", size=11, bold=True, color=RED)

label(s, 0.9, 6.3, W - 1.8, 0.5,
      "§6.2: when a pre_X interception blocks, post_X is NOT emitted. "
      "§6.3: any failure (interceptor raised, timeout, invalid verdict) "
      "fails closed to deny with a host_error:* reason.",
      size=10, color=SLATE)

footer(s, 6, TOTAL)
notes(s, "evaluate_only mode (§8): same pipeline, verdicts recorded, action "
         "always proceeds. The host MUST NOT present an evaluate_only "
         "outcome as enforcement.")


# ---------------------------------------------------------------------------
# 7. Escalate / approval sequence diagram
s = prs.slides.add_slide(BLANK)
header(s, "SPEC §9", "The escalate approval seam",
       "Approval is bound to context_identity so the approver consents to "
       "exactly the action that will execute")

actors = [("Host", BLUE), ("Interceptor", MAGENTA), ("Resolver", AMBER),
          ("Action", GREEN)]
ax = {}
top, life_bot = 2.2, 6.7
aw = (W - 1.8) / len(actors)
for i, (a, col) in enumerate(actors):
    x = 0.9 + i * aw + aw / 2
    ax[a] = x
    label(s, x - 0.8, top, 1.6, 0.45, a, fill=col, color=CARD,
          size=11, bold=True)
    connector(s, x, top + 0.45, x, life_bot, color=RULE, width=1.0,
              arrow=False, dashed=True)

def msg(y, frm, to, text, *, col=SLATE, ret=False):
    connector(s, ax[frm], y, ax[to], y, color=col, width=1.5,
              dashed=ret)
    mid = (ax[frm] + ax[to]) / 2
    label(s, mid - 1.6, y - 0.32, 3.2, 0.28, text, size=8, color=col,
          font=MONO)

y = 3.1
msg(y, "Host", "Interceptor", "intercept(ctx)  id=sha256(L0+L1)", col=BLUE); y += 0.45
msg(y, "Interceptor", "Host", "Verdict{escalate}", col=MAGENTA, ret=True); y += 0.45
msg(y, "Host", "Resolver", "ApprovalRequest{id, ctx, verdict}", col=AMBER); y += 0.45
msg(y, "Resolver", "Host", "ApprovalResolution{outcome, id', verdict'}",
    col=AMBER, ret=True); y += 0.45

# decision diamond
dy = y + 0.05
label(s, ax["Host"] - 0.95, dy, 1.9, 0.5, "id' == id ?",
      fill=CARD, line=NAVY, size=9, bold=True, font=MONO)
y += 0.7
msg(y, "Host", "Action", "approve→permit  →  execute(target)", col=GREEN); y += 0.4
label(s, ax["Host"] + 0.3, y - 0.12, 4.5, 0.3,
      "reject / unresolved / id mismatch  →  deny (host_error:approval_*)",
      size=8, color=RED, font=MONO, align=PP_ALIGN.LEFT)

footer(s, 7, TOTAL)
notes(s, "Open design question RM-09: should ApprovalResolution.verdict be "
         "restricted to allow|warn so a resolver cannot inject a transform "
         "that bypasses the interceptor chain?")


# ---------------------------------------------------------------------------
# 8. Canonical identity
s = prs.slides.add_slide(BLANK)
header(s, "SPEC §10", "context_identity: stable across SDKs, scoped to L0+L1",
       "Approval binding and cross-host audit correlation depend on "
       "byte-identical hashes")

# left: what's hashed
hx, hw = 1.0, 5.4
label(s, hx, 2.1, hw, 0.4, "canonical_json( strip_to_L0+L1(ctx) )",
      fill=CODE_BG, color=CODE_FG, size=11, font=MONO)
inc = [("spec, interception_point, timestamp, sequence", True),
       ("agent.{id, framework}", True),
       ("session.{id}", True),
       ("target  +  L1 fields for this interception_point", True),
       ("agent.{name, version}, session.started_at", False),
       ("trace, tenant, budgets, usage", False),
       ("extensions.*", False)]
for i, (t, included) in enumerate(inc):
    col = GREEN if included else SUBTLE
    mark = "✓" if included else "✗"
    label(s, hx, 2.65 + i * 0.45, 0.4, 0.4, mark, fill=None, color=col,
          size=14, bold=True)
    label(s, hx + 0.45, 2.65 + i * 0.45, hw - 0.45, 0.4, t,
          align=PP_ALIGN.LEFT, size=10, color=GRAPHITE if included else SUBTLE,
          font=MONO)

connector(s, hx + hw + 0.2, 3.8, hx + hw + 1.0, 3.8, color=BLUE, width=2.5)
label(s, hx + hw + 1.1, 3.5, 1.5, 0.6, "SHA-256",
      fill=BLUE, color=CARD, size=12, bold=True)
connector(s, hx + hw + 2.7, 3.8, hx + hw + 3.5, 3.8, color=BLUE, width=2.5)
label(s, hx + hw + 3.6, 3.55, 2.4, 0.5, "sha256:7c3f…", fill=CODE_BG,
      color=CYAN, size=11, font=MONO)

tf = textbox(s, hx + hw + 1.1, 4.6, W - hx - hw - 1.8, 1.8)
txt(tf, "Why strip L2/L3?", size=11, bold=True, color=NAVY)
txt(tf, "Adding a trace_id or budget MUST NOT change the approval hash.",
    size=10, color=GRAPHITE, first=False)
txt(tf, "input_identity vs enforced_identity:", size=11, bold=True,
    color=NAVY, first=False)
txt(tf, "Equal except after a transform; the approver consents to the "
        "enforced one.", size=10, color=GRAPHITE, first=False)

label(s, 0.9, 6.45, W - 1.8, 0.45,
      "RM-02/03 (high): canonical JSON underspecified vs RFC 8785; "
      "no cross-SDK golden vectors yet → identities may diverge silently",
      size=10, color=AMBER, fill=None, line=AMBER, align=PP_ALIGN.LEFT)

footer(s, 8, TOTAL)
notes(s, "Open question: replace §10.1 with a normative reference to "
         "RFC 8785 (JCS) or document the deltas explicitly.")


# ---------------------------------------------------------------------------
# 9. CTK architecture + conformance levels
s = prs.slides.add_slide(BLANK)
header(s, "SPEC §13", "Conformance Test Kit",
       "Language-agnostic JSON vectors; one Harness shim per framework; "
       "three conformance levels")

# left: pyramid
px, pw = 1.0, 4.6
levels_p = [
    ("Level 3  Complete", "L2 fields, result_labels propagation, "
     "parallel/streaming, stable identity", CYAN),
    ("Level 2  Enforcing", "honours deny / transform / escalate / "
     "evaluate_only", BLUE),
    ("Level 1  Instrumented", "all interception points fire in order with valid L0+L1 "
     "context", SLATE),
]
for i, (t, d, col) in enumerate(levels_p):
    inset = (2 - i) * 0.45
    sp = box(s, px + inset, 2.2 + i * 1.25, pw - 2 * inset, 1.1,
             fill=col, line=None)
    tf = sp.text_frame
    tf.word_wrap = True
    txt(tf, t, size=12, bold=True, color=CARD, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(tf, d, size=8, color=RGBColor(0xE8, 0xE8, 0xE8),
        align=PP_ALIGN.CENTER, first=False)

# right: CTK pipeline
rx = px + pw + 0.7
steps = [
    ("vectors/*.json", "scenario + interceptor_script + expect", CARD, MAGENTA),
    ("CTK runner", "ScriptedInterceptor + RecordingInterceptor", CARD, BLUE),
    ("Harness", "framework shim: setup() / run() / teardown()",
     RGBColor(0xFF, 0xF6, 0xE5), AMBER),
    ("framework under test", "real adapter, mock model + tools", CARD, GREEN),
]
for i, (t, d, fill, line) in enumerate(steps):
    y = 2.2 + i * 1.05
    sp = label(s, rx, y, W - rx - 0.9, 0.85, "", fill=fill, line=line)
    tf = sp.text_frame
    txt(tf, t, size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(tf, d, size=9, color=SLATE, align=PP_ALIGN.LEFT, first=False)
    if i < len(steps) - 1:
        connector(s, rx + 0.4, y + 0.85, rx + 0.4, y + 1.05,
                  color=SLATE, width=1.5)

label(s, rx, 6.45, W - rx - 0.9, 0.4,
      "14 vectors today (L1+L2). RM-01 (critical): negative-path / "
      "fail-closed vectors missing.", size=9, color=RED,
      align=PP_ALIGN.LEFT)

footer(s, 9, TOTAL)
notes(s, "Harness is the ONLY framework-specific code (~100 lines). "
         "Vectors and runner are shared. A claim is "
         "(framework, adapter-ver, spec-ver, Level N, sdk-lang@ver).")


# ---------------------------------------------------------------------------
# 10. Multi-lang status matrix
s = prs.slides.add_slide(BLANK)
header(s, "STATUS", "SDK parity at v0.1.0-alpha.1",
       "Python is the reference; others ship types + identity + Harness "
       "interface; CTK runners scaffolded")

cols = ["types", "canonical /\nidentity", "emitter", "Harness\niface",
        "CTK\nrunner", "CI"]
rows = [
    ("Python", [2, 2, 2, 2, 2, 2]),
    ("Rust", [2, 2, 0, 2, 1, 2]),
    ("TypeScript", [2, 2, 0, 2, 1, 2]),
    (".NET", [2, 2, 0, 2, 1, 2]),
    ("Go", [2, 2, 0, 2, 1, 2]),
]
cell_col = {2: GREEN, 1: AMBER, 0: RULE}
cell_txt = {2: "done", 1: "scaffold", 0: "—"}

mx, my = 1.6, 2.4
cw, ch = 1.55, 0.7
for j, c in enumerate(cols):
    label(s, mx + 1.6 + j * cw, my, cw - 0.05, ch, c,
          size=9, bold=True, color=NAVY)
for i, (lang, vals) in enumerate(rows):
    y = my + ch + 0.1 + i * (ch + 0.08)
    label(s, mx, y, 1.55, ch, lang, size=11, bold=True, color=NAVY,
          align=PP_ALIGN.LEFT)
    for j, v in enumerate(vals):
        label(s, mx + 1.6 + j * cw, y, cw - 0.05, ch, cell_txt[v],
              fill=cell_col[v], color=CARD if v else SUBTLE,
              size=9, bold=True)

tf = textbox(s, 0.9, 6.4, W - 1.8, 0.6)
txt(tf, "47/47 Python tests pass. All 14 vectors pass against "
        "ReferenceHarness. 6/6 CI jobs green at HEAD.",
    size=10, color=GRAPHITE)

footer(s, 10, TOTAL)
notes(s, "Emitter (host-side §6-§9 helper) only exists in Python; other "
         "SDKs expose types + identity + interfaces and expect the host "
         "framework to wire them. RM-27 tracks CTK runner parity.")


# ---------------------------------------------------------------------------
# 11. Findings: severity x effort scatter
s = prs.slides.add_slide(BLANK)
header(s, "ARCHITECTURAL REVIEW", "58 deduplicated findings: severity vs effort",
       "18-dimension parallel review; 96 raw findings deduplicated to 58")

# axes
ox, oy, aw, ah = 1.6, 6.4, 8.4, 3.9
connector(s, ox, oy, ox + aw, oy, color=GRAPHITE, width=1.5)  # x
connector(s, ox, oy, ox, oy - ah, color=GRAPHITE, width=1.5, arrow=False)  # y
for i, e in enumerate(["S", "M", "L", "XL"]):
    label(s, ox + (i + 0.5) * aw / 4 - 0.3, oy + 0.05, 0.6, 0.3, e,
          size=10, color=SLATE)
label(s, ox + aw / 2 - 1.0, oy + 0.4, 2.0, 0.3, "effort →", size=10,
      color=SLATE)
for i, sv in enumerate(["low", "medium", "high", "critical"]):
    label(s, ox - 1.0, oy - (i + 0.5) * ah / 4 - 0.15, 0.95, 0.3, sv,
          size=10, color=SLATE, align=PP_ALIGN.RIGHT)
label(s, ox - 1.4, oy - ah / 2 - 0.15, 0.3, 0.3, "↑", size=12, color=SLATE)

# "Now" zone
box(s, ox + 0.05, oy - ah, aw / 2 - 0.05, ah / 2, fill=None,
    line=GREEN, line_w=2.0)
label(s, ox + 0.1, oy - ah + 0.05, 1.4, 0.3, "NOW (26)",
      size=10, bold=True, color=GREEN, align=PP_ALIGN.LEFT)

# bubbles (severity 1-4, effort 1-4, count)
buckets = {
    (4, 2, "RM-01"): 1,
    (3, 1, ""): 12, (3, 2, ""): 13, (3, 3, "RM-27"): 1,
    (2, 1, ""): 17, (2, 2, ""): 8,
    (1, 1, ""): 5, (1, 2, ""): 1,
}
sev_col = {4: RED, 3: AMBER, 2: BLUE, 1: SUBTLE}
for (sv, ef, tag), n in buckets.items():
    cx_b = ox + (ef - 0.5) * aw / 4
    cy_b = oy - (sv - 0.5) * ah / 4
    r = 0.12 + 0.06 * math.sqrt(n)
    box(s, cx_b - r, cy_b - r, 2 * r, 2 * r, fill=sev_col[sv],
        shape=MSO_SHAPE.OVAL, line=None)
    label(s, cx_b - r, cy_b - 0.15, 2 * r, 0.3, str(n),
          size=10, bold=True, color=CARD)
    if tag:
        label(s, cx_b + r + 0.05, cy_b - 0.12, 0.9, 0.25, tag,
              size=8, color=sev_col[sv], font=MONO, align=PP_ALIGN.LEFT)

# right: top findings
tf = textbox(s, ox + aw + 0.5, 2.4, W - ox - aw - 1.2, 4.0)
txt(tf, "Top of NOW", size=12, bold=True, color=NAVY)
for fid, t in [
    ("RM-01", "Fail-closed has zero CTK coverage"),
    ("RM-08", "host_error:* verdicts violate own schema"),
    ("RM-16", "emit() result is ignorable; deny bypassable"),
    ("RM-04", "Mutable ctx → covert transform"),
    ("RM-02/03", "Canonical JSON diverges; no golden vectors"),
    ("RM-05", "No threat model / non-goals"),
    ("RM-19/20", "Actions unpinned; Cargo.lock ignored"),
]:
    txt(tf, f"  {fid}", size=9, bold=True, color=AMBER, font=MONO,
        first=False, space_after=0)
    txt(tf, f"    {t}", size=9, color=GRAPHITE, first=False)

footer(s, 11, TOTAL)
notes(s, "Full findings: docs/arch-review/2026-06-22/ROADMAP.md. "
         "1 critical, 40 high, 46 medium, 9 low before dedup.")


# ---------------------------------------------------------------------------
# 12. Open design questions
s = prs.slides.add_slide(BLANK)
header(s, "DECISIONS NEEDED", "Six open design questions",
       "Each changes normative spec text; need a call before v0.1.0")

qs = [
    ("Canonical JSON", "Replace §10.1 with RFC 8785 (JCS) reference, or "
     "document deltas as a JCS profile?", "RM-02"),
    ("Approval transform", "May ApprovalResolution return transform, or "
     "restrict to allow|warn so resolvers cannot bypass interceptors?", "RM-09"),
    ("Zero interceptors", "Implicit allow, or deny with "
     "host_error:no_consumer?", "RM-12"),
    ("Multi-interceptor combine", "Keep first-block short-circuit, or invoke "
     "all and apply deny>escalate>transform>warn>allow lattice?", "RM-11"),
    ("Interceptor trust roles", "All interceptors one trust class, or define "
     "observing vs controlling at registration?", "RM-29"),
    ("Streaming egress", "In scope for 0.1 (buffer-until-permit), or "
     "explicit non-goal deferred to 0.2?", "RM-10"),
]
for i, (t, q, ref) in enumerate(qs):
    y = 2.1 + i * 0.78
    label(s, 0.9, y, 0.5, 0.65, str(i + 1), fill=NAVY, color=CARD,
          size=14, bold=True)
    sp = box(s, 1.5, y, W - 2.4, 0.65, fill=CARD, line=RULE)
    tf = sp.text_frame
    tf.word_wrap = True
    txt(tf, f"{t}", size=11, bold=True, color=NAVY,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)
    txt(tf, f"{q}   [{ref}]", size=9, color=GRAPHITE,
        align=PP_ALIGN.LEFT, first=False)

footer(s, 12, TOTAL)
notes(s, "These are genuine design choices, not bugs. Each has a defensible "
         "answer either way; the spec must pick one and state it normatively.")


# ---------------------------------------------------------------------------
# 13. Roadmap swim lanes
s = prs.slides.add_slide(BLANK)
header(s, "ROADMAP", "Now / Next / Later",
       "docs/arch-review/2026-06-22/ROADMAP.md")

lanes = [
    ("NOW", "before widening collaborator access", GREEN, 26,
     ["RM-01 fail-closed CTK vectors",
      "RM-02/03 RFC 8785 + golden vectors",
      "RM-04/16 immutable ctx + must-use result",
      "RM-05 threat model + non-goals",
      "RM-08 schema fix for host_error:*",
      "RM-09/11/12 §7 + §9 spec tightening",
      "RM-19/20/25 CI pinning + CODEOWNERS"]),
    ("NEXT", "before spec/v0.1.0 tag", BLUE, 26,
     ["RM-27 cross-SDK CTK parity gate",
      "RM-29 interceptor trust roles",
      "RM-34 per-interceptor verdict provenance",
      "RM-40 payload size/depth bounds",
      "RM-43 release workflow + SBOM + signing",
      "RM-45 Level-3 vectors",
      "RM-51 OWASP/NIST controls map"]),
    ("LATER", "post-0.1", SLATE, 6,
     ["RM-54 pluggable telemetry sink",
      "RM-55 per-SDK READMEs",
      "RM-56/57/58 hygiene"]),
]
lw = (W - 1.8 - 0.4) / 3
for i, (t, sub, col, n, items) in enumerate(lanes):
    x = 0.9 + i * (lw + 0.2)
    sp = box(s, x, 2.0, lw, 0.7, fill=col, line=None)
    tf = sp.text_frame
    txt(tf, f"{t}  ({n})", size=14, bold=True, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(tf, sub, size=8, color=RGBColor(0xE8, 0xE8, 0xE8),
        align=PP_ALIGN.CENTER, first=False)
    box(s, x, 2.7, lw, 4.0, fill=CARD, line=col)
    tf2 = textbox(s, x + 0.15, 2.85, lw - 0.3, 3.7)
    for j, it in enumerate(items):
        txt(tf2, it, size=9, color=GRAPHITE, first=(j == 0), space_after=6)

footer(s, 13, TOTAL)
notes(s, "Now horizon is critical/high with S/M effort. Full list with "
         "evidence and recommendations in ROADMAP.md.")


# ---------------------------------------------------------------------------
# 14. Ask
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
tf = textbox(s, 1.2, 1.6, W - 2.4, 4.5)
txt(tf, "ASK", size=12, bold=True, color=CYAN)
txt(tf, "What we need from this review", size=26, bold=True,
    color=RGBColor(0xFF, 0xFF, 0xFF), first=False, space_after=14)
asks = [
    ("Decide the six open questions", "Each is a one-paragraph spec edit "
     "once decided; blocking RM-02/09/10/11/12/29."),
    ("Confirm the layering cut", "Verdict and approval seam in agent-hooks "
     "vs ACS; identity computed over AgentContext not PolicyInput."),
    ("Nominate per-language CTK owners", "Rust, TS, .NET, Go runners are "
     "scaffolded; need an owner each to reach parity (RM-27)."),
    ("Approve the conformance-claim format", "CLAIMS.md row shape and the "
     "scope disclaimer (RM-30)."),
]
for t, d in asks:
    txt(tf, f"  •  {t}", size=14, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF), first=False, space_after=2)
    txt(tf, f"      {d}", size=10, color=RGBColor(0xC0, 0xC8, 0xD4),
        first=False, space_after=10)

label(s, 1.2, H - 0.9, W - 2.4, 0.4,
      "github.com/responsibleai/agent-hooks  ·  spec/AGENT-HOOKS-0.1.md  ·  "
      "docs/arch-review/2026-06-22/",
      size=9, color=RGBColor(0x90, 0x9C, 0xAA), font=MONO,
      align=PP_ALIGN.LEFT)

footer(s, 14, TOTAL)
notes(s, "")


# ---------------------------------------------------------------------------
out = "agent-hooks-design-review.pptx"
prs.save(out)
print(f"wrote {out} ({len(prs.slides)} slides)")
